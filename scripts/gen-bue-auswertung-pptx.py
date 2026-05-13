"""
Generate Auswertung-BUe-Stichwoche.pptx: 1 Slide mit Funnel + Phasen + Top-10-Pools.
Layout: 16:9 (13.33 x 7.5 inch).

Datenquelle: Mappe10.xlsx (Stichwoche April 2026, n=1.094 analysierte Mails).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Farbpalette - ERGO-Naehe
ERGO_RED    = RGBColor(0xC8, 0x10, 0x2E)
ERGO_DARK   = RGBColor(0x1F, 0x29, 0x37)
PHASE1_BG   = RGBColor(0x16, 0x65, 0x34)   # gruen - Auto-Kandidat
PHASE2_BG   = RGBColor(0xCA, 0x8A, 0x04)   # gelb - Klaerung
PHASE3_BG   = RGBColor(0x6B, 0x72, 0x80)   # grau - out-of-scope
GREY_BG     = RGBColor(0xE5, 0xE7, 0xEB)
GREY_BORDER = RGBColor(0x9C, 0xA3, 0xAF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1F, 0x29, 0x37)
ACCENT_BLUE = RGBColor(0x1E, 0x40, 0xAF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

def add_textbox(left, top, width, height, text, font_size=10, bold=False,
                color=DARK_TEXT, fill=None, border=None, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border:
        tb.line.color.rgb = border
        tb.line.width = Pt(0.75)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2);  tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_multiline(left, top, width, height, lines, font_size=10, font_bold=False,
                  fill=None, border=None, color=DARK_TEXT, align=PP_ALIGN.LEFT,
                  line_spacing=1.05):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border:
        tb.line.color.rgb = border; tb.line.width = Pt(0.75)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(3);  tf.margin_bottom = Pt(3)
    tf.word_wrap = True
    for i, item in enumerate(lines):
        opts = {}
        if isinstance(item, tuple):
            txt, opts = item
        else:
            txt = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", "Calibri")
        r.font.size = Pt(opts.get("size", font_size))
        r.font.bold = opts.get("bold", font_bold)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb

# Title bar
title_h = Inches(0.55)
title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, title_h)
title_bar.fill.solid(); title_bar.fill.fore_color.rgb = ERGO_DARK
title_bar.line.fill.background()
title_bar.text_frame.margin_left = Inches(0.4)
title_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = title_bar.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "Automatisierungs-Potenzial Bestandsübertragung ERGO   |   Stichwoche April 2026   |   1.094 Mails analysiert"
r.font.name = "Calibri"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

# Kernaussage banner (under title)
kern_top = title_h + Inches(0.08)
kern_h = Inches(0.45)
kern = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), kern_top, prs.slide_width - Inches(0.6), kern_h)
kern.fill.solid(); kern.fill.fore_color.rgb = PHASE1_BG
kern.line.fill.background()
kern.text_frame.margin_left = Inches(0.2); kern.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = kern.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "▶  Kernaussage: 59 % der Mails sind regelbasiert automatisierbare BÜ-Erstvorgänge mit vollständiger MV  →  646 von 1.094 Mails in einer einzigen Stichwoche."
r.font.name = "Calibri"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

# ---------- FUNNEL (LINKS) ----------
funnel_top  = kern_top + kern_h + Inches(0.18)
funnel_left = Inches(0.3)
funnel_w    = Inches(4.7)

add_textbox(funnel_left, funnel_top, funnel_w, Inches(0.32),
            "FUNNEL: Vom Eingang zum Automatisierungs-Kandidaten",
            font_size=12, bold=True, color=ERGO_DARK)

# Funnel-Stufen (echte Zahlen aus Mappe10.xlsx)
stages = [
    ("1.094",  "Eingang gesamt",                              "100 %", WHITE, ERGO_DARK),
    ("– 64",   "Rauschen (Bounce-NDR, System, ERGO-Outbound)", "",      GREY_BG, DARK_TEXT),
    ("1.030",  "Echte Makler-Vorgänge",                       "94 %",  WHITE, ERGO_DARK),
    ("– 36",   "Andere Klassifikation (Antrag, Anfrage)",     "",      GREY_BG, DARK_TEXT),
    ("994",    "BÜ-Vorgänge",                                  "91 %",  WHITE, ERGO_DARK),
    ("– 103",  "Reminder / Wiedervorlage",                    "",      GREY_BG, DARK_TEXT),
    ("891",    "BÜ-Erstvorgänge",                              "81 %",  WHITE, ERGO_DARK),
    ("– 245",  "MV-Klärungsbedarf",                           "",      GREY_BG, DARK_TEXT),
    ("646",    "AUTOMATISIERUNGS-KANDIDATEN",                  "59 %",  PHASE1_BG, WHITE),
]

y = funnel_top + Inches(0.32)
row_h = Inches(0.36)
for absol, label, pct, bg, fg in stages:
    is_loss = absol.startswith("–")
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, funnel_left, y, funnel_w, row_h)
    box.fill.solid(); box.fill.fore_color.rgb = bg
    if is_loss:
        box.line.color.rgb = GREY_BORDER; box.line.width = Pt(0.5)
    else:
        box.line.color.rgb = ERGO_DARK if fg == WHITE else GREY_BORDER
        box.line.width = Pt(0.5)
    box.text_frame.margin_left = Pt(8); box.text_frame.margin_right = Pt(8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = absol + "   "
    r1.font.name = "Calibri"; r1.font.size = Pt(11 if not is_loss else 9)
    r1.font.bold = True; r1.font.italic = is_loss; r1.font.color.rgb = fg
    r2 = p.add_run(); r2.text = label
    r2.font.name = "Calibri"; r2.font.size = Pt(10 if not is_loss else 9)
    r2.font.italic = is_loss; r2.font.color.rgb = fg
    if pct:
        r3 = p.add_run(); r3.text = "    " + pct
        r3.font.name = "Calibri"; r3.font.size = Pt(11)
        r3.font.bold = True; r3.font.color.rgb = fg
    y = y + row_h + Inches(0.02)

y_end = y + Inches(0.04)
add_multiline(funnel_left, y_end, funnel_w, Inches(0.4),
              [("Definition Automatisierungs-Kandidat: BÜ-Vorgang + Erstvorgang (kein Reminder) + MV-Status OK.",
                {"size": 8, "italic": True, "color": RGBColor(0x6b, 0x72, 0x80)})])

# ---------- MITTLERE SPALTE: Phasen ----------
mid_left = funnel_left + funnel_w + Inches(0.25)
mid_w    = Inches(4.0)

add_textbox(mid_left, funnel_top, mid_w, Inches(0.32),
            "Drei Phasen der Bearbeitung",
            font_size=12, bold=True, color=ERGO_DARK)

# Phase 1 Box
p1_top = funnel_top + Inches(0.34)
p1_h   = Inches(1.65)
add_multiline(mid_left, p1_top, mid_w, p1_h, [
    ("PHASE 1   Sofort automatisierbar   59 %  (646)", {"size": 11, "bold": True, "color": WHITE}),
    ("Sparten: Komposit 53 %  |  KV 29 %  |  Leben 16 %", {"size": 9, "color": WHITE}),
    ("Geschäftstyp (BÜ): einfacher Vertrag 76 %  |  Kundenverbindung 24 %", {"size": 9, "color": WHITE}),
    ("Anhang-Qualität: 94 % Text-PDF (automatisch parsebar)", {"size": 9, "color": WHITE}),
    ("Pool-Konzentration: Top 3 Pools decken 54 % ab", {"size": 9, "color": WHITE, "bold": True}),
    ("→ Schnittstellen-Priorisierung auf JDC, Fonds Finanz, blau direkt", {"size": 9, "color": WHITE, "italic": True}),
], fill=PHASE1_BG, border=PHASE1_BG)

# Phase 2 Box
p2_top = p1_top + p1_h + Inches(0.12)
p2_h   = Inches(1.55)
add_multiline(mid_left, p2_top, mid_w, p2_h, [
    ("PHASE 2   Klärungsbedarf   22 %  (245)", {"size": 11, "bold": True, "color": WHITE}),
    ("35 %   MV fehlt im Anhang         →  Outbound: MV nachfordern", {"size": 9, "color": WHITE}),
    ("33 %   MV unvollständig            →  Klärungs-Template", {"size": 9, "color": WHITE}),
    ("29 %   kein Anhang                  →  Outbound: MV nachfordern", {"size": 9, "color": WHITE}),
    ("4 %    nicht prüfbar (Scan/OCR)   →  manuelle Sichtung", {"size": 9, "color": WHITE}),
    ("→ Halbautomatischer Workflow mit Standard-Anschreiben", {"size": 9, "color": WHITE, "italic": True}),
], fill=PHASE2_BG, border=PHASE2_BG)

# Phase 3 Box
p3_top = p2_top + p2_h + Inches(0.12)
p3_h   = Inches(1.35)
add_multiline(mid_left, p3_top, mid_w, p3_h, [
    ("Außerhalb Scope   19 %  (203)", {"size": 11, "bold": True, "color": WHITE}),
    ("9 %   Reminder / Wiedervorlage   →  eigener SLA-Use-Case", {"size": 9, "color": WHITE}),
    ("3 %   Nicht-BÜ (Anfrage, Antrag) →  andere Pipeline", {"size": 9, "color": WHITE}),
    ("6 %   Rauschen                    →  Spam-Filter / NDR", {"size": 9, "color": WHITE}),
    ("5 %   Sondertarif / Flotte         →  Konsortien manuell", {"size": 9, "color": WHITE}),
], fill=PHASE3_BG, border=PHASE3_BG)

# ---------- RECHTE SPALTE: Top Pools ----------
right_left = mid_left + mid_w + Inches(0.25)
right_w    = prs.slide_width - right_left - Inches(0.3)

add_textbox(right_left, funnel_top, right_w, Inches(0.32),
            "Top 10 Pools in Automatisierungs-Kandidaten (n = 646)",
            font_size=12, bold=True, color=ERGO_DARK)

# echte Zahlen aus Mappe10.xlsx (BÜe + Erstvorgang + MV OK)
pools = [
    ("1.  JDC / Jung, DMS & Cie.",                     33, True),
    ("2.  Fonds Finanz Maklerservice",                 12, True),
    ("3.  blau direkt",                                10, True),
    ("4.  DEMA Deutsche Versicherungsmakler",           5, False),
    ("5.  Fuhrmann Versicherungsmakler",                3, False),
    ("6.  Qualitypool",                                 3, False),
    ("7.  Securess Versicherungsmakler",                2, False),
    ("8.  [pma:] Finanz-Service",                       2, False),
    ("9.  Apella AG",                                   2, False),
    ("10. TauRes Gesellschaft für Investmentberatung",  2, False),
]

pool_y = funnel_top + Inches(0.34)
pool_row_h = Inches(0.32)
pool_label_w = Inches(2.4)
pool_bar_max_w = right_w - pool_label_w - Inches(0.55)
max_pct = 33
for label, pct, is_top3 in pools:
    lbl = slide.shapes.add_textbox(right_left, pool_y, pool_label_w, pool_row_h)
    lbl.fill.background(); lbl.line.fill.background()
    lbl.text_frame.margin_left = Pt(2); lbl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = lbl.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
    rr = pp.add_run(); rr.text = label
    rr.font.name = "Calibri"; rr.font.size = Pt(9.5)
    rr.font.bold = is_top3; rr.font.color.rgb = DARK_TEXT
    bar_bg_left = right_left + pool_label_w
    bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_bg_left,
                                     pool_y + Inches(0.07), pool_bar_max_w, Inches(0.18))
    bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = GREY_BG
    bg_bar.line.fill.background()
    val_w = Emu(int(pool_bar_max_w * pct / max_pct))
    if val_w > 0:
        val_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_bg_left,
                                          pool_y + Inches(0.07), val_w, Inches(0.18))
        val_bar.fill.solid()
        val_bar.fill.fore_color.rgb = ACCENT_BLUE if is_top3 else RGBColor(0x60, 0xa5, 0xfa)
        val_bar.line.fill.background()
    pct_lbl = slide.shapes.add_textbox(bar_bg_left + pool_bar_max_w + Inches(0.05),
                                        pool_y, Inches(0.55), pool_row_h)
    pct_lbl.fill.background(); pct_lbl.line.fill.background()
    pct_lbl.text_frame.margin_left = Pt(2); pct_lbl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp2 = pct_lbl.text_frame.paragraphs[0]; pp2.alignment = PP_ALIGN.LEFT
    rr2 = pp2.add_run(); rr2.text = str(pct) + " %"
    rr2.font.name = "Calibri"; rr2.font.size = Pt(9.5)
    rr2.font.bold = is_top3; rr2.font.color.rgb = ACCENT_BLUE if is_top3 else DARK_TEXT
    pool_y = pool_y + pool_row_h

sum_y = pool_y + Inches(0.05)
add_multiline(right_left, sum_y, right_w, Inches(0.4), [
    ("Top 3 = 54 %   |   Top 10 = 74 %   |   Long Tail (90+ Pools) = 26 %",
     {"size": 9, "italic": True, "color": RGBColor(0x4b, 0x55, 0x63)})
])

# ---------- UNTEN: Jahres-Hochrechnung + Empfehlung ----------
bottom_top = Inches(6.05)
bottom_h   = Inches(1.4)

# Jahres-Hochrechnung (1.094 Mails/Woche × ca. 52 Wochen)
hr_w = Inches(6.5)
add_multiline(Inches(0.3), bottom_top, hr_w, bottom_h, [
    ("Hochrechnung auf das Jahr (1.094 Mails × ca. 52 Wochen)", {"size": 11, "bold": True, "color": ERGO_DARK}),
    ("ca. 33.600   Phase 1   Automatisierungs-Potenzial", {"size": 10, "color": DARK_TEXT, "bold": True}),
    ("ca. 12.700   Phase 2   Halbautomatische MV-Klärung", {"size": 10, "color": DARK_TEXT}),
    ("ca.  5.400   Phase 3   Reminder- / SLA-Logik", {"size": 10, "color": DARK_TEXT}),
    ("ca.  1.900             Andere Klassifikation (Antrag, Anfrage)", {"size": 10, "color": DARK_TEXT}),
    ("ca.  3.300             Rauschen (NDR, System, Outbound)", {"size": 10, "color": DARK_TEXT}),
    ("ca. 56.900   GESAMT-Mails p.a. im SHUK-Innendienst (Hochrechnung)", {"size": 10, "color": ERGO_DARK, "bold": True}),
], fill=RGBColor(0xF3, 0xF4, 0xF6), border=GREY_BORDER)

# Empfehlung rechts
emp_left = Inches(0.3) + hr_w + Inches(0.25)
emp_w = prs.slide_width - emp_left - Inches(0.3)
add_multiline(emp_left, bottom_top, emp_w, bottom_h, [
    ("Empfehlung", {"size": 11, "bold": True, "color": WHITE}),
    ("Phase-1-Pilot mit Top 3 Pools (JDC, Fonds Finanz, blau direkt) deckt mit nur DREI Schnittstellen rund 54 % des Automatisierungs-Potenzials ab — Top 10 decken 74 %.",
     {"size": 11, "color": WHITE}),
    ("", {"size": 4}),
    ("Folge-Phase: MV-Klärungs-Templates für die drei Hauptgründe (MV fehlt, MV unvollständig, kein Anhang) → +22 % Volumen halbautomatisch.",
     {"size": 10, "color": WHITE, "italic": True}),
], fill=ERGO_RED, border=ERGO_RED)

# Footer / Datenquelle
add_textbox(Inches(0.3), Inches(7.25), Inches(13.0), Inches(0.22),
            "Datenquelle: Mappe10.xlsx — ergo-vorgang-analyse v2.13  |  KI-Klassifikation Mail + PDF-Anhang  |  Pool-Schreibweisen konsolidiert  |  Eine Stichwoche, keine Saisonalität berücksichtigt",
            font_size=7.5, italic=True, color=RGBColor(0x6b, 0x72, 0x80))

out = "/home/user/emailimport/Auswertung-BUe-Stichwoche-Vorschlag.pptx"
prs.save(out)
import os
print("Saved:", out, "size:", os.path.getsize(out), "bytes")
