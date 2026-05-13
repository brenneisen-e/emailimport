"""
Generate Prozess-IST-SOLL-Vorschlag.pptx: 1 Slide mit Vorher/Nachher-Prozess.
16:9 Layout, IST oben, SOLL unten, Mehrwert rechts hervorgehoben.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Farb-Palette
ERGO_DARK  = RGBColor(0x1F, 0x29, 0x37)
ERGO_RED   = RGBColor(0xC8, 0x10, 0x2E)
IST_BG     = RGBColor(0xFE, 0xE2, 0xE2)   # helles Rosa - IST
IST_BORDER = RGBColor(0xDC, 0x26, 0x26)   # Rot
SOLL_BG    = RGBColor(0xD1, 0xFA, 0xE5)   # helles Gruen - SOLL
SOLL_BORDER= RGBColor(0x05, 0x96, 0x69)   # Gruen
PROBLEM_BG = RGBColor(0xFE, 0xCA, 0xCA)
PROBLEM_BORDER = RGBColor(0xB9, 0x1C, 0x1C)
WIN_BG     = RGBColor(0xA7, 0xF3, 0xD0)
WIN_BORDER = RGBColor(0x04, 0x7B, 0x55)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TEXT  = RGBColor(0x4B, 0x55, 0x63)
GREY_BG    = RGBColor(0xF3, 0xF4, 0xF6)
GREY_BORDER= RGBColor(0x9C, 0xA3, 0xAF)
STEP_BG    = RGBColor(0xE0, 0xE7, 0xFF)   # hellblau fuer Prozessschritte
STEP_BORDER= RGBColor(0x4F, 0x46, 0xE5)
MAKRO_BG   = RGBColor(0xFE, 0xF3, 0xC7)   # gelb - der entscheidende Schritt
MAKRO_BORDER = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def add_run(p, text, *, size=11, bold=False, italic=False, color=ERGO_DARK, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r

def textbox(left, top, w, h, lines, *, fill=None, border=None, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE, margin=4):
    tb = slide.shapes.add_textbox(left, top, w, h)
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border:
        tb.line.color.rgb = border; tb.line.width = Pt(0.75)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_left = Pt(margin); tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin); tf.margin_bottom = Pt(margin)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        opts = {}
        if isinstance(item, tuple):
            txt, opts = item
        else:
            txt = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("ls", 1.1)
        add_run(p, txt,
                size=opts.get("size", 10),
                bold=opts.get("bold", False),
                italic=opts.get("italic", False),
                color=opts.get("color", ERGO_DARK))
    return tb

def shape_box(left, top, w, h, text, *, fill, border, font_size=10, bold=False, color=ERGO_DARK):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1.0)
    sh.text_frame.margin_left = Pt(6); sh.text_frame.margin_right = Pt(6)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    sh.text_frame.word_wrap = True
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    add_run(p, text, size=font_size, bold=bold, color=color)
    return sh

def arrow(x1, y1, x2, y2, color=ERGO_DARK, width=2.5):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(width)
    # Arrow head
    line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    # Try arrow head via OOXML
    tail = ln.find(qn('a:headEnd'))
    if tail is None:
        from lxml import etree
        head = etree.SubElement(ln, qn('a:tailEnd'))
        head.set('type', 'triangle')
        head.set('w', 'med')
        head.set('len', 'med')
    return line

# ---------- TITLE ----------
title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
title_bar.fill.solid(); title_bar.fill.fore_color.rgb = ERGO_DARK
title_bar.line.fill.background()
title_bar.text_frame.margin_left = Inches(0.4)
title_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = title_bar.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
add_run(p, "Bestandsuebertragung ERGO   |   IST-Prozess  →  SOLL-Prozess mit Makro-Vorqualifizierung",
        size=16, bold=True, color=WHITE)

# ---------- IST-BLOCK (oben) ----------
ist_top = Inches(0.75)
ist_h = Inches(2.7)

# Linke Section: IST-Header
ist_label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), ist_top, Inches(1.1), Inches(0.5))
ist_label.fill.solid(); ist_label.fill.fore_color.rgb = IST_BORDER
ist_label.line.fill.background()
ist_label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = ist_label.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, "IST", size=18, bold=True, color=WHITE)

# IST-Subtitle
textbox(Inches(1.5), ist_top, Inches(7), Inches(0.5),
        [("Email kommt rein → wird blind durchgereicht, keine inhaltliche Pruefung",
          {"size": 11, "italic": True, "color": GREY_TEXT})],
        anchor=MSO_ANCHOR.MIDDLE)

# IST-Prozesskette
step_y = ist_top + Inches(0.65)
step_h = Inches(1.0)
step_w = Inches(2.1)
gap = Inches(0.25)
ist_start_x = Inches(0.3)

ist_steps = [
    ("Email + Anhang", "vom Maklerpool"),
    ("Weiterleitung", "blind an KWAD"),
    ("Indizierung", "rein technisch"),
    ("Postkorb-\nEinsortierung", "Manuelle Sichtung"),
]

x = ist_start_x
for i, (label, sub) in enumerate(ist_steps):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, step_y, step_w, step_h)
    sh.fill.solid(); sh.fill.fore_color.rgb = STEP_BG
    sh.line.color.rgb = STEP_BORDER; sh.line.width = Pt(1.0)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    sh.text_frame.word_wrap = True
    sh.text_frame.margin_left = Pt(6); sh.text_frame.margin_right = Pt(6)
    p1 = sh.text_frame.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    add_run(p1, label, size=11, bold=True, color=ERGO_DARK)
    p2 = sh.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    add_run(p2, sub, size=8, italic=True, color=GREY_TEXT)
    # Arrow zur naechsten Box
    if i < len(ist_steps) - 1:
        arrow(x + step_w, step_y + step_h // 2,
              x + step_w + gap, step_y + step_h // 2,
              color=ERGO_DARK)
    x = x + step_w + gap

# IST-Probleme rechts
ist_prob_left = x + Inches(0.05)
ist_prob_w = prs.slide_width - ist_prob_left - Inches(0.3)
textbox(ist_prob_left, step_y, ist_prob_w, step_h,
        [
            ("PROBLEME im IST-Prozess", {"size": 11, "bold": True, "color": PROBLEM_BORDER}),
            ("•  teilweise gar keine BUe (Antrag/Anfrage/Schaden im BUe-Postkorb)", {"size": 9.5, "color": ERGO_DARK}),
            ("•  teilweise nur Reminder (kein neuer Vorgang)", {"size": 9.5, "color": ERGO_DARK}),
            ("•  keine Vorqualifizierung (Pool / Sparte / Vertragsnummer / MV-Status)", {"size": 9.5, "color": ERGO_DARK}),
            ("→ Aufwand entsteht erst im Postkorb beim Sachbearbeiter", {"size": 9.5, "italic": True, "color": PROBLEM_BORDER}),
        ],
        fill=PROBLEM_BG, border=PROBLEM_BORDER, anchor=MSO_ANCHOR.TOP, margin=8)

# ---------- TRENNER ----------
sep_y = ist_top + ist_h + Inches(0.15)
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), sep_y, prs.slide_width - Inches(0.6), Inches(0.02))
sep.fill.solid(); sep.fill.fore_color.rgb = GREY_BORDER
sep.line.fill.background()

# ---------- SOLL-BLOCK (unten) ----------
soll_top = sep_y + Inches(0.18)
soll_h = Inches(3.4)

# SOLL-Header
soll_label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), soll_top, Inches(1.1), Inches(0.5))
soll_label.fill.solid(); soll_label.fill.fore_color.rgb = SOLL_BORDER
soll_label.line.fill.background()
soll_label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = soll_label.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, "SOLL", size=18, bold=True, color=WHITE)

# SOLL-Subtitle
textbox(Inches(1.5), soll_top, Inches(8.5), Inches(0.5),
        [("Email → Makro nutzt ErgoGPT zur Vorqualifizierung → strukturierte Notiz wird mitgeliefert",
          {"size": 11, "italic": True, "color": GREY_TEXT})],
        anchor=MSO_ANCHOR.MIDDLE)

# SOLL-Prozesskette
step_y2 = soll_top + Inches(0.65)
step_h2 = Inches(1.2)
step_w2 = Inches(2.0)
gap2 = Inches(0.18)

soll_steps = [
    ("Email +\nAnhang", "vom Maklerpool", STEP_BG, STEP_BORDER, ERGO_DARK),
    ("MAKRO\nErgoGPT-Analyse", "Vorqualifizierung", MAKRO_BG, MAKRO_BORDER, ERGO_DARK),
    ("Strukturierte\nNotiz erzeugen", "als zusaetzl. Anhang", STEP_BG, STEP_BORDER, ERGO_DARK),
    ("Betreff\nanpassen", "klassifiziert", STEP_BG, STEP_BORDER, ERGO_DARK),
    ("Bessere\nIndizierung", "auf KI-Metadaten", STEP_BG, STEP_BORDER, ERGO_DARK),
    ("Postkorb +\nNotiz als Image", "im Vorgang sichtbar", STEP_BG, STEP_BORDER, ERGO_DARK),
]

x = Inches(0.3)
for i, (label, sub, fill, border, color) in enumerate(soll_steps):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, step_y2, step_w2, step_h2)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1.5 if i == 1 else 1.0)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    sh.text_frame.word_wrap = True
    sh.text_frame.margin_left = Pt(4); sh.text_frame.margin_right = Pt(4)
    p1 = sh.text_frame.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    add_run(p1, label, size=10.5, bold=True, color=color)
    p2 = sh.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    add_run(p2, sub, size=8, italic=True, color=GREY_TEXT)
    if i < len(soll_steps) - 1:
        arrow(x + step_w2, step_y2 + step_h2 // 2,
              x + step_w2 + gap2, step_y2 + step_h2 // 2,
              color=ERGO_DARK)
    x = x + step_w2 + gap2

# SOLL-Mehrwerte unter den Boxen
mw_top = step_y2 + step_h2 + Inches(0.15)
mw_h   = Inches(1.25)

# Linke Box: Was die Notiz enthaelt
textbox(Inches(0.3), mw_top, Inches(6.2), mw_h,
        [
            ("Strukturierte Notiz enthaelt:", {"size": 10, "bold": True, "color": SOLL_BORDER}),
            ("• Klassifikation (BUe / Anfrage / Antrag / Schaden / Nicht-Standard)", {"size": 9, "color": ERGO_DARK}),
            ("• Vorgangstyp (Makler-Vorgang / Bounce / System / Outbound)", {"size": 9, "color": ERGO_DARK}),
            ("• Erstvorgang vs. Reminder (+ Quelle: Anhang-PDF / Betreff / Body)", {"size": 9, "color": ERGO_DARK}),
            ("• Maklerpool (konsolidiert) + Maklernummer + Sparte + Vertragsnummer", {"size": 9, "color": ERGO_DARK}),
            ("• MV-Status (OK / unvollstaendig / fehlt) + Schluesselworte", {"size": 9, "color": ERGO_DARK}),
            ("• Mail-Zusammenfassung (1-2 Saetze, neutral)", {"size": 9, "color": ERGO_DARK}),
        ],
        fill=SOLL_BG, border=SOLL_BORDER, anchor=MSO_ANCHOR.TOP, margin=8)

# Rechte Box: Mehrwerte fuer den Sachbearbeiter
textbox(Inches(6.7), mw_top, Inches(6.4), mw_h,
        [
            ("Mehrwerte fuer den Sachbearbeiter:", {"size": 10, "bold": True, "color": WIN_BORDER}),
            ("✓  Reminder werden gleich aussortiert (kein Doppel-Aufwand)", {"size": 9, "color": ERGO_DARK}),
            ("✓  Nicht-BUe wandert in passenden Postkorb (keine Falsch-Einsortierung)", {"size": 9, "color": ERGO_DARK}),
            ("✓  Vorqualifikation sofort sichtbar (Pool, MV, Sparte)", {"size": 9, "color": ERGO_DARK}),
            ("✓  Volltext-Suche im Postkorb auf KI-Metadaten moeglich", {"size": 9, "color": ERGO_DARK}),
            ("✓  Auswertung / Reporting auf realer Datenbasis", {"size": 9, "color": ERGO_DARK}),
            ("✓  Voraussetzung fuer Voll-Automatisierung (Phase 2)", {"size": 9, "color": WIN_BORDER, "bold": True}),
        ],
        fill=WIN_BG, border=WIN_BORDER, anchor=MSO_ANCHOR.TOP, margin=8)

# Footer
textbox(Inches(0.3), Inches(7.25), Inches(13), Inches(0.22),
        [("Datenbasis Wirksamkeitsnachweis: 5.000 Mails / Stichwoche, KI-Klassifikation Mail + PDF | siehe Auswertungs-Slide",
          {"size": 7.5, "italic": True, "color": GREY_TEXT})])

out = "/home/user/emailimport/Prozess-IST-SOLL-Vorschlag.pptx"
prs.save(out)
import os
print("Saved:", out, "size:", os.path.getsize(out), "bytes")
